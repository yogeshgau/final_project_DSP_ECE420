"""
Generate ECE420 project presentation using python-pptx.
Run: py ml/make_slides.py
Output: ECE420_Project_Slides.pptx in project root
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT_PATH = os.path.join(os.path.dirname(__file__), '..', 'ECE420_Project_Slides.pptx')

# ── Color palette ────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT     = RGBColor(0x1E, 0x88, 0xE5)   # bright blue
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xB0, 0xBE, 0xC5)
GOLD       = RGBColor(0xFF, 0xD5, 0x4F)

# ── Slide content ────────────────────────────────────────────────────────────
SLIDES = [
    # (title, subtitle_or_none, [bullets])
    # bullets can be strings (level-1) or (indent_level, text) tuples

    # 1 ── TITLE SLIDE
    {
        "title": "IMU-Based Exercise Recognition\nand Repetition Counter",
        "subtitle": "ECE420 — Embedded Signal Processing\nAndroid Application with Real-Time ML Inference",
        "bullets": [],
        "is_title_slide": True,
    },

    # 2 ── SYSTEM OVERVIEW
    {
        "title": "System Overview",
        "subtitle": "End-to-End Pipeline",
        "bullets": [
            "Goal: classify exercise type and count repetitions in real time using only the phone's built-in IMU",
            "Three target exercises: Biceps Press (BP), Shoulder Press (SP), Triceps Extension (TP)",
            "Pipeline spans hardware sensing → signal processing → ML training → on-device inference → UI",
            "",
            "Four major subsystems:",
            (1, "Sensor layer — Android SensorManager, accelerometer + gyroscope at ~50 Hz"),
            (1, "Signal processing — IIR low-pass filter, jerk magnitude, autocorrelation, peak detection"),
            (1, "Machine-learning pipeline — offline Python training (sklearn), logistic regression exported to Java"),
            (1, "Android application — real-time classification, rep counting, calorie estimation, session history"),
            "",
            "All inference runs entirely on-device; no network calls or cloud dependency",
        ],
    },

    # 3 ── DATA COLLECTION
    {
        "title": "Data Collection Infrastructure",
        "subtitle": "Android IMU Recording Module",
        "bullets": [
            "Custom ImuCollector class wraps Android SensorManager and buffers raw sensor events",
            "Two physical sensors registered simultaneously:",
            (1, "TYPE_ACCELEROMETER — linear acceleration in m/s²  (ax, ay, az)"),
            (1, "TYPE_GYROSCOPE    — angular velocity in rad/s     (gx, gy, gz)"),
            "",
            "Sampling rate: SENSOR_DELAY_GAME (~50 Hz nominal)",
            "Each event stores event.timestamp in nanoseconds for precise timing reconstruction",
            "",
            "Recording session exports a 7-column CSV per set:",
            (1, "timestamp_ms,  ax,  ay,  az,  gx,  gy,  gz"),
            "",
            "Export uses ACTION_CREATE_DOCUMENT (Storage Access Framework) — no WRITE_EXTERNAL_STORAGE permission required",
            "Files collected offline; labeled by folder: data/bp/, data/sp/, data/Triceps/",
        ],
    },

    # 4 ── SENSOR CONFIGURATION
    {
        "title": "Sensor Configuration & Sampling",
        "subtitle": "Hardware Layer Details",
        "bullets": [
            "SENSOR_DELAY_GAME is the closest Android constant to 50 Hz (~20 ms event interval)",
            "Actual rate varies by device; timestamp is used rather than assuming fixed Δt",
            "",
            "During active recording all six channels buffered in List<Float> per axis:",
            (1, "accelX, accelY, accelZ  — one value per SensorEvent from TYPE_ACCELEROMETER"),
            (1, "gyroX,  gyroY,  gyroZ   — one value per SensorEvent from TYPE_GYROSCOPE"),
            (1, "accelTs, gyroTs          — List<Long> of nanosecond timestamps for each sensor"),
            "",
            "Sensor fusion not used; accelerometer and gyroscope are treated as independent channels",
            "Gravity component is NOT removed — raw accelerometer values are used directly",
            (1, "Gravity provides a stable DC offset that encodes phone orientation, useful for classification"),
            "",
            "Three parallel ImuCollector instances run simultaneously:",
            (1, "imuCollector   — dedicated to rep counting signal"),
            (1, "dataCollector  — dedicated to CSV file export"),
            (1, "liveClassifier — feeds the real-time ML inference module"),
        ],
    },

    # 5 ── SIGNAL PRE-PROCESSING: LPF
    {
        "title": "Signal Pre-Processing: Low-Pass Filter",
        "subtitle": "IIR First-Order Butterworth Approximation",
        "bullets": [
            "Raw accelerometer data contains high-frequency noise from phone vibration and sensor quantization",
            "A causal first-order IIR filter is applied independently to each of the six axes:",
            (1, "y[n]  =  α · x[n]  +  (1 − α) · y[n−1]"),
            "",
            "Filter coefficient α = 0.76, chosen to give a −3 dB cutoff of approximately 14 Hz:",
            (1, "ωc  =  arccos( 1 − α² / (2(1−α)) )"),
            (1, "At fs = 50 Hz: fc = ωc · fs / (2π)  ≈  14 Hz"),
            "",
            "Why 14 Hz:",
            (1, "Exercise repetition rates fall between 0.3 and 2.5 Hz (fundamental + harmonics up to ~10 Hz)"),
            (1, "Cutting at 14 Hz preserves all biomechanically relevant harmonics"),
            (1, "Previous value of α = 0.2 (~1.8 Hz cutoff) was too aggressive and distorted rep waveforms"),
            "",
            "Filter is applied causally (left-to-right), so no look-ahead latency — suitable for streaming",
            "Initial condition: y[0] = x[0]  (avoids startup transient)",
        ],
    },

    # 6 ── JERK MAGNITUDE
    {
        "title": "Jerk Magnitude Computation",
        "subtitle": "Deriving a Periodic Scalar Signal from 3-Axis Data",
        "bullets": [
            "After filtering, each axis is differentiated discretely to obtain the jerk (rate of change of acceleration):",
            (1, "Δax[n] = ax[n+1] − ax[n]"),
            (1, "Δay[n] = ay[n+1] − ay[n]"),
            (1, "Δaz[n] = az[n+1] − az[n]"),
            "",
            "The scalar jerk magnitude collapses all three axes into a single non-negative signal:",
            (1, "J[n]  =  √( Δax² + Δay² + Δaz² )"),
            "",
            "Why differentiate before taking magnitude:",
            (1, "Differentiation suppresses the gravity DC offset, making the signal zero-mean around rest"),
            (1, "Each repetition produces two distinct jerk peaks (concentric + eccentric phase)"),
            (1, "The magnitude operation makes the signal rotation-invariant — result is the same regardless of phone orientation"),
            "",
            "Output length is N − 1 samples (one shorter than the input due to differencing)",
            "Same pipeline applied to gyroscope channels as an alternative rep-counting signal",
        ],
    },

    # 7 ── AUTOCORRELATION
    {
        "title": "Repetition Period Estimation via Autocorrelation",
        "subtitle": "Finding the Dominant Periodicity in the Jerk Signal",
        "bullets": [
            "Repetitive exercise produces a quasi-periodic jerk signal; the period corresponds to one full rep",
            "Autocorrelation R(τ) measures self-similarity at lag τ:",
            (1, "R(τ)  =  (1/T) · Σ_t  s(t) · s(t+τ)"),
            (1, "where s is the zero-meaned jerk magnitude and T is the number of valid pairs at lag τ"),
            "",
            "Zero-meaning is critical: jerk magnitude is always non-negative (DC offset).",
            "Without zero-meaning, R(τ) is dominated by the constant offset and the true periodic peak is buried",
            "",
            "Search range for lag τ:",
            (1, "Minimum:  20 samples  ≈ 0.4 s  (fastest plausible rep rate)"),
            (1, "Maximum: 250 samples  ≈ 5.0 s  (slowest plausible rep rate)"),
            (1, "Upper limit also capped at N/2 to ensure statistical reliability"),
            "",
            "The lag τ* with the highest R(τ) is taken as the estimated repetition period L",
            "L is passed directly to the peak detection stage as the expected rep interval",
        ],
    },

    # 8 ── PEAK DETECTION
    {
        "title": "Peak Detection Algorithm",
        "subtitle": "Counting Reps from Jerk Peaks",
        "bullets": [
            "Peaks in the jerk magnitude correspond to the high-effort phases of each repetition",
            "",
            "Three simultaneous criteria must be satisfied for a sample i to be accepted as a peak:",
            "",
            "1. Amplitude threshold:",
            (1, "J[i]  ≥  mean(J)  +  0.4 · std(J)"),
            (1, "Rejects low-energy bumps and baseline noise while remaining robust to signal amplitude variation"),
            "",
            "2. Local maximum:",
            (1, "J[i]  >  J[j]  for all j in [i − L/2, i + L/2], j ≠ i"),
            (1, "Window half-width = L/2 ensures only the tallest peak within one half-period is selected"),
            "",
            "3. Minimum separation:",
            (1, "i − lastAcceptedPeak  ≥  L/2  samples"),
            (1, "Prevents double-counting a single rep due to a broad or noisy jerk peak"),
            "",
            "Final rep count = number of accepted peaks",
            "Peak indices also used for inter-peak timing to compute average rep cadence (s/rep)",
        ],
    },

    # 9 ── REP COUNTING SUMMARY
    {
        "title": "Repetition Counting Pipeline Summary",
        "subtitle": "Complete Signal Chain",
        "bullets": [
            "Stage 1 — Raw IMU: ax, ay, az at ~50 Hz  (List<Float>, updated on each SensorEvent)",
            "",
            "Stage 2 — IIR Low-Pass Filter (α = 0.76, fc ≈ 14 Hz):",
            (1, "Applied independently to ax, ay, az → smoothed signals fa, fb, fc"),
            "",
            "Stage 3 — Jerk Magnitude:",
            (1, "Discrete derivative of each filtered axis → L2 norm → scalar J[n]"),
            "",
            "Stage 4 — Autocorrelation Period Estimation:",
            (1, "Zero-mean J, compute R(τ) for τ ∈ [20, 250], return τ* = argmax R(τ)"),
            "",
            "Stage 5 — Peak Detection:",
            (1, "Find local maxima in J with minimum separation τ*/2, above mean + 0.4·std"),
            (1, "Rep count = number of accepted peaks"),
            "",
            "Stage 6 — Cadence (s/rep):",
            (1, "avgPeakIntervalSec = mean of consecutive peak-to-peak gaps / 50 Hz"),
            "",
            "Identical pipeline available for gyroscope (countRepsGyro); accelerometer used by default",
        ],
    },

    # 10 ── FEATURE EXTRACTION OVERVIEW
    {
        "title": "Feature Extraction: Sliding Window Approach",
        "subtitle": "Bridging Signal Processing and Machine Learning",
        "bullets": [
            "Classification requires a fixed-length numerical descriptor of each exercise window",
            "",
            "Window parameters:",
            (1, "WINDOW = 100 samples  ≈ 2 seconds of motion at 50 Hz"),
            (1, "STEP   =  50 samples  →  50% overlap between consecutive windows"),
            (1, "50% overlap doubles the number of training examples per session"),
            "",
            "Each window of 100 samples × 6 channels is reduced to a 46-dimensional feature vector",
            "",
            "Three feature families computed per window:",
            (1, "Per-channel statistics  (36 features): 6 stats × 6 channels"),
            (1, "Magnitude statistics    ( 4 features): accel magnitude + gyro magnitude, each with mean & std"),
            (1, "Pearson cross-correlations ( 6 features): 6 axis pairs across both sensors"),
            "",
            "Feature computation is deterministic and runs on the device in Java, matching the Python training code exactly",
            "All features are scale-normalized (z-score) before classification using the training-set scaler",
        ],
    },

    # 11 ── PER-CHANNEL STATS
    {
        "title": "Per-Channel Statistical Features",
        "subtitle": "36 Features — 6 Statistics × 6 Channels",
        "bullets": [
            "Six statistics are computed from each channel's 100-sample window:",
            "",
            (1, "Mean        — captures the gravitational and postural DC component"),
            (1, "Std         — measures the intensity/energy of movement"),
            (1, "Min         — lower bound of the motion arc"),
            (1, "Max         — upper bound of the motion arc"),
            (1, "Range       — Max − Min; total amplitude sweep of the movement"),
            (1, "RMS         — √(mean(x²)); sensitive to both signal level and variance"),
            "",
            "Applied to all six channels in order: ax, ay, az, gx, gy, gz",
            "This produces 6 × 6 = 36 features occupying indices [0 … 35] of the feature vector",
            "",
            "Why include both Std and Range:",
            (1, "Std is sensitive to outliers; Range captures the true excursion; both complement each other"),
            "",
            "Why include RMS alongside Mean:",
            (1, "For zero-mean signals Std = RMS; for non-zero-mean signals (gravity-biased axes) they diverge,"),
            (1, "and RMS captures signal power more accurately than Std alone"),
        ],
    },

    # 12 ── MAGNITUDE + PEARSON
    {
        "title": "Global Features: Magnitude & Cross-Correlation",
        "subtitle": "10 Features Encoding Multi-Axis Relationships",
        "bullets": [
            "Magnitude features (indices 36–39) — collapse 3D motion to a scalar signal:",
            (1, "Accel magnitude:  ||a||[n] = √( ax² + ay² + az² )   → mean and std over window"),
            (1, "Gyro  magnitude:  ||g||[n] = √( gx² + gy² + gz² )   → mean and std over window"),
            (1, "Magnitude is orientation-invariant: same value regardless of phone mount direction"),
            "",
            "Pearson cross-correlations (indices 40–45) — encode exercise-specific coupling between axes:",
            (1, "r(ax, ay),  r(ax, az),  r(ay, az)  — accelerometer inter-axis correlations"),
            (1, "r(gx, gy),  r(gx, gz),  r(gy, gz)  — gyroscope  inter-axis correlations"),
            "",
            "r(a, b)  =  Σ(aᵢ−ā)(bᵢ−b̄)  /  √[ Σ(aᵢ−ā)² · Σ(bᵢ−b̄)² ]",
            "",
            "Biological interpretation:",
            (1, "Biceps press has strong az–ay coupling (elbow flexion in the sagittal plane)"),
            (1, "Shoulder press has strong ay–az coupling dominated by vertical axis rotation"),
            (1, "Triceps extension has a distinct gx–gz gyro signature from elbow extension direction"),
            "",
            "Pearson r is bounded [−1, 1] and scale-invariant — robust across subjects of different strength",
        ],
    },

    # 13 ── FULL FEATURE VECTOR
    {
        "title": "Full 46-Dimensional Feature Vector",
        "subtitle": "Index Map",
        "bullets": [
            "Indices  0 –  5:   ax  { mean, std, min, max, range, rms }",
            "Indices  6 – 11:   ay  { mean, std, min, max, range, rms }",
            "Indices 12 – 17:   az  { mean, std, min, max, range, rms }",
            "Indices 18 – 23:   gx  { mean, std, min, max, range, rms }",
            "Indices 24 – 29:   gy  { mean, std, min, max, range, rms }",
            "Indices 30 – 35:   gz  { mean, std, min, max, range, rms }",
            "Index   36:        accel magnitude — mean",
            "Index   37:        accel magnitude — std",
            "Index   38:        gyro  magnitude — mean",
            "Index   39:        gyro  magnitude — std",
            "Index   40:        Pearson r( ax, ay )",
            "Index   41:        Pearson r( ax, az )",
            "Index   42:        Pearson r( ay, az )",
            "Index   43:        Pearson r( gx, gy )",
            "Index   44:        Pearson r( gx, gz )",
            "Index   45:        Pearson r( gy, gz )",
            "",
            "Feature vector is identical in Python (train.py) and Java (ExerciseClassifier.java) — any mismatch causes silent accuracy loss",
        ],
    },

    # 14 ── DATASET
    {
        "title": "Dataset Construction & Labeling",
        "subtitle": "Session-Based IMU Recordings",
        "bullets": [
            "Data collected using the app's built-in CSV recorder across multiple sessions:",
            (1, "Biceps Press  (label 0, folder data/bp/)   — 6 session files"),
            (1, "Shoulder Press (label 1, folder data/sp/)   — 6 session files"),
            (1, "Triceps Extension (label 2, folder data/Triceps/) — 4 session files"),
            "",
            "Each session file: one continuous recording of 10–40+ repetitions",
            "Session files vary in length; shorter sessions produce fewer windows",
            "",
            "Sliding window extraction per session:",
            (1, "WINDOW = 100, STEP = 50  →  floor((N − 100)/50) + 1 windows per file"),
            (1, "Each window independently labeled with the session's exercise class"),
            "",
            "Total dataset: ~400–600 feature vectors depending on recording lengths",
            "",
            "Missing data handling:",
            (1, "np.nan_to_num(X, nan=0.0) applied before training"),
            (1, "NaN can arise in Pearson r when a channel is constant across the window"),
            "",
            "No data augmentation — diversity comes from natural variation across sessions",
        ],
    },

    # 15 ── CROSS-VALIDATION
    {
        "title": "Cross-Validation Strategy",
        "subtitle": "GroupKFold — Session-Based Holdout",
        "bullets": [
            "Standard k-fold cross-validation would cause data leakage in time-series window data:",
            (1, "Adjacent windows from the same session share almost identical samples"),
            (1, "A window from session 3 in training and an overlapping window from session 3 in test"),
            (1, "would make the test accuracy artifically high and not generalize to new users"),
            "",
            "Solution: GroupKFold with groups = session ID",
            (1, "Each fold holds out exactly one complete session as the test set"),
            (1, "All windows from that session are withheld from training"),
            (1, "No window from a test session ever appears in training — zero leakage"),
            "",
            "Number of folds = number of unique sessions (16 total: 6+6+4)",
            "Each fold trains on 15 sessions, tests on 1 — leave-one-session-out cross-validation",
            "",
            "cross_val_predict aggregates out-of-fold predictions across all folds",
            "Final reported accuracy/F1 is across all windows, each predicted under a held-out-session regime",
            "",
            "After cross-validation, models are re-fit on the full dataset for deployment",
        ],
    },

    # 16 ── CLASSIFIERS
    {
        "title": "Classifier Selection: LR and Random Forest",
        "subtitle": "Training & Evaluation",
        "bullets": [
            "Two classifiers evaluated during development:",
            "",
            "Logistic Regression (multinomial, one-vs-rest disabled):",
            (1, "Pipeline: StandardScaler → LogisticRegression(multi_class='multinomial', C=1.0, max_iter=2000)"),
            (1, "StandardScaler: z-normalizes each feature to zero mean and unit variance"),
            (1, "Multinomial softmax loss over all three classes simultaneously"),
            (1, "Regularization strength C=1.0 (L2 penalty) prevents overfitting on small dataset"),
            (1, "Closed-form scaler parameters (mean, scale) and LR coefficients (coef, bias) exported to JSON"),
            "",
            "Random Forest:",
            (1, "RandomForestClassifier(n_estimators=200, max_depth=None)"),
            (1, "Serves as a non-linear upper bound on achievable accuracy with these features"),
            (1, "Feature importances reported — identifies which of the 46 features matter most"),
            (1, "Not deployable to Android (no efficient Java path for tree ensembles)"),
            "",
            "Decision: Logistic Regression chosen for on-device deployment",
            (1, "Forward pass is a single matrix multiply + softmax — O(46 × 3) = 138 multiply-adds"),
            (1, "Trivially implementable in Java without any ML library"),
        ],
    },

    # 17 ── ON-DEVICE INFERENCE
    {
        "title": "On-Device Inference Pipeline",
        "subtitle": "Exporting and Running the Classifier in Java",
        "bullets": [
            "Weight export chain (Python side):",
            (1, "train.py trains LR pipeline, writes weights.json: { bias, mean, scale, coef }"),
            (1, "export_weights.py reads weights.json, generates ExerciseClassifier.java with literal float arrays"),
            (1, "Generated file replaces the placeholder version in the Android source tree"),
            "",
            "On-device forward pass (Java side):",
            (1, "Step 1 — Extract 46 features from the current 100-sample window"),
            (1, "Step 2 — Z-score normalize: z[i] = (feat[i] − SCALER_MEAN[i]) / SCALER_SCALE[i]"),
            (1, "Step 3 — Compute logit per class k: logit[k] = BIAS[k] + COEF[k] · z"),
            (1, "Step 4 — Numerically stable softmax: subtract max logit before exponentiation"),
            (1, "Step 5 — Predicted class = argmax(probs);  confidence = max(probs)"),
            "",
            "Numerical stability of softmax:",
            (1, "Computing exp(logit[k]) directly can overflow for large logits"),
            (1, "Subtracting max(logits) before exponentiation is algebraically equivalent and stays in float range"),
            "",
            "No ML framework (TFLite, ONNX) needed — pure Java arithmetic, zero extra APK size",
        ],
    },

    # 18 ── LIVE CLASSIFICATION & GATE
    {
        "title": "Real-Time Classification & Idle Detection Gate",
        "subtitle": "LiveClassifier — Streaming Inference Every ~3 Seconds",
        "bullets": [
            "LiveClassifier runs alongside ImuCollector, processing the rolling sensor buffer:",
            (1, "STEP = 150 samples — classification fires approximately every 3 seconds at 50 Hz"),
            (1, "Each inference operates on the most recent 100-sample window (2 seconds)"),
            "",
            "Two-stage gate before calling the classifier — both must pass:",
            "",
            "Gate 1 — Motion energy threshold:",
            (1, "Compute accel magnitude std over the window: σ = std(||a||)"),
            (1, "MOTION_THRESHOLD = 0.5 m/s²  — rejects stationary or very slow movement"),
            (1, "If σ < 0.5: emit null result → UI shows 'Rest', calories = 0"),
            "",
            "Gate 2 — Periodicity check (rep counter used as detector):",
            (1, "Run SignalProcessor.countRepsAccel() on the same 100-sample window"),
            (1, "If rep count < 1: emit null result → UI shows 'Rest', calories = 0"),
            (1, "Ensures random aperiodic movement (walking, adjusting equipment) is not classified as exercise"),
            "",
            "Only if both gates pass: run ExerciseClassifier.classify() and emit the Result",
            "Inference runs on a background thread (AsyncTask) to avoid blocking the UI thread",
        ],
    },

    # 19 ── CALORIE ESTIMATION
    {
        "title": "Calorie Estimation Model",
        "subtitle": "MET-Based Metabolic + Mechanical Work",
        "bullets": [
            "Two additive components estimate total energy expenditure per session:",
            "",
            "Component 1 — Metabolic cost (MET-based):",
            (1, "Calories = MET × body_weight_kg × duration_hours"),
            (1, "MET values: Biceps Press = 4.0,  Shoulder Press = 5.0,  Triceps Extension = 3.5"),
            (1, "Body weight entered by user (kg), persisted across sessions via SharedPreferences"),
            (1, "Duration = number of confirmed live-classified windows × window_duration"),
            "",
            "Component 2 — Mechanical work (load-based):",
            (1, "Calories = (load_kg × g × ROM_m × 2 × reps) / (muscle_efficiency × J_per_cal)"),
            (1, "load_kg: weight the user reports carrying (dumbbell / barbell)"),
            (1, "ROM (range of motion): BP = 0.35 m,  SP = 0.55 m,  TP = 0.40 m"),
            (1, "Factor of 2: accounts for both concentric and eccentric phases"),
            (1, "Muscle efficiency ≈ 25%;  1 dietary calorie = 4184 J"),
            "",
            "Display:",
            (1, "Live panel shows rolling calorie total while exercise is ongoing"),
            (1, "During rest/idle: live calories display 0"),
            (1, "Session total accumulates across all sets and is shown in the stats bar"),
        ],
    },

    # 20 ── SESSION MANAGEMENT & UI
    {
        "title": "Session Management & User Interface",
        "subtitle": "Multi-Set Workflow",
        "bullets": [
            "Session lifecycle:",
            (1, "Start button: begins simultaneous IMU collection across all three collectors"),
            (1, "Next Set button: freezes current set, processes and logs it, immediately restarts for next set"),
            (1, "Stop button: ends collection, runs full rep-counting and stats pipeline, displays results"),
            "",
            "Set history panel (dynamically populated ListView-style):",
            (1, "Each set entry shows: set number, exercise label, rep count, duration, calories"),
            (1, "Color-coded by exercise: Biceps Press = blue, Shoulder Press = green, Triceps Extension = orange"),
            (1, "Clear button resets session history"),
            "",
            "Stats bar (always visible):",
            (1, "Total session duration, average rep cadence (s/rep), peak acceleration (m/s²), cumulative calories"),
            "",
            "Live detection panel (updates every ~3 s):",
            (1, "Exercise label in class color, softmax confidence percentage, rolling live calorie display"),
            (1, "Shows 'Rest  /  0 kcal' when either motion or periodicity gate rejects the window"),
            "",
            "IMU recording panel (independent of rep counting):",
            (1, "Record / Stop buttons save raw CSV at any time via the system file picker"),
            (1, "Used during data collection for ML training; unchanged during deployment"),
        ],
    },

    # 21 ── PERFORMANCE METRICS
    {
        "title": "Performance Statistics Displayed to User",
        "subtitle": "On-Device Signal Analysis Results",
        "bullets": [
            "Rep count:",
            (1, "Derived from peak detection on the jerk magnitude — reported per set and cumulatively"),
            "",
            "Average rep cadence (s/rep):",
            (1, "avgPeakIntervalSec = mean of peak-to-peak gaps / 50 Hz"),
            (1, "More accurate than (total_time / rep_count) because it measures the actual rhythmic interval"),
            (1, "Computed from the same autocorrelation → peak pipeline used for counting"),
            "",
            "Peak acceleration (m/s²):",
            (1, "Maximum accel magnitude sample over the entire set"),
            (1, "Reflects the explosive force at the fastest point of the lift"),
            "",
            "Session duration:",
            (1, "Wall-clock time from Start to Stop, reported in MM:SS format"),
            "",
            "All statistics are computed on a background thread after Stop/Next Set is pressed",
            "UI thread receives results via a Handler post — avoids ANR on long signal buffers",
        ],
    },

    # 22 ── THREE-CLASS EXTENSION
    {
        "title": "Three-Class Extension: Adding Triceps Extension",
        "subtitle": "Scaling from Binary to Multinomial Classification",
        "bullets": [
            "Initial system: binary classification (Biceps Press vs. Shoulder Press)",
            "Extension: added Triceps Extension as a third class (label 2)",
            "",
            "Changes required in the ML pipeline:",
            (1, "FOLDERS updated: added entry 2 → 'Triceps' pointing to data/Triceps/"),
            (1, "LogisticRegression: multi_class='multinomial' already set — no structural change"),
            (1, "BIAS vector grows from 2 elements to 3"),
            (1, "COEF matrix grows from 2×46 to 3×46"),
            (1, "weights.json and ExerciseClassifier.java regenerated by the export script"),
            "",
            "Changes required in Android:",
            (1, "LABEL_TP = 2, LABEL_NAMES[2] = 'Triceps Extension' added to ExerciseClassifier"),
            (1, "labelColor(): TP → deep orange (#E65100)"),
            (1, "labelMet():   TP → MET 3.5"),
            (1, "getRomM():    TP → ROM 0.40 m"),
            (1, "All switch/if chains in RepCounterActivity updated to handle label 0, 1, or 2"),
            "",
            "The softmax inference loop is class-count agnostic — loop over N_CLASSES requires no change",
            "Feature extraction is unchanged — exercise-discriminating information already present in the 46-dim vector",
        ],
    },

    # 23 ── KEY DESIGN DECISIONS
    {
        "title": "Key Design Decisions & Tradeoffs",
        "subtitle": "Engineering Choices Made During Development",
        "bullets": [
            "IIR over FIR for LPF:",
            (1, "IIR uses one coefficient and one state variable — O(1) per sample, minimal memory"),
            (1, "FIR would require storing ~50+ past samples for equivalent roll-off"),
            "",
            "Jerk magnitude over raw magnitude:",
            (1, "Raw accel magnitude is dominated by gravity (~9.8 m/s²); DC bias obscures exercise signal"),
            (1, "Differencing removes DC and emphasizes dynamic transitions that mark rep boundaries"),
            "",
            "Autocorrelation over fixed threshold for period:",
            (1, "Rep rate varies per user and fatigue level; fixed threshold misses slow or fast reps"),
            (1, "Autocorrelation adapts the detection window to the actual cadence observed in the data"),
            "",
            "Logistic Regression over Random Forest for deployment:",
            (1, "LR: 3×46 = 138 floats, one matrix multiply — fits in phone L1 cache"),
            (1, "RF: 200 trees × ~1000 nodes each — ~MB of memory + cache-unfriendly traversal"),
            "",
            "GroupKFold over standard k-fold:",
            (1, "Window overlap between consecutive windows of the same session causes leakage"),
            (1, "Session-level holdout gives an honest estimate of generalization to new recording sessions"),
        ],
    },
]

# ── PPTX builder ─────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def hex_rgb(h):
    return RGBColor((h >> 16) & 0xFF, (h >> 8) & 0xFF, h & 0xFF)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def title_para(tf, text, size=36, bold=True, color=WHITE):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_accent_bar(slide):
    """Thin horizontal bar below the title area."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.5), Inches(1.75),
        Inches(12.33), Pt(3)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def build_title_slide(prs, data):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    set_bg(slide, BG_DARK)

    # Main title
    tb = add_textbox(slide, Inches(0.6), Inches(1.8), Inches(12.1), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = data["title"]
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Accent bar
    bar = slide.shapes.add_shape(1, Inches(0.6), Inches(4.3), Inches(4.0), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # Subtitle
    tb2 = add_textbox(slide, Inches(0.6), Inches(4.5), Inches(12.1), Inches(1.6))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = data["subtitle"]
    run2.font.size = Pt(22)
    run2.font.bold = False
    run2.font.color.rgb = LIGHT_GREY


def build_content_slide(prs, data):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    set_bg(slide, BG_DARK)

    # Title box
    tb_title = add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12.33), Inches(0.75))
    tf_title = tb_title.text_frame
    p = tf_title.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = data["title"]
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Subtitle / section label
    if data.get("subtitle"):
        tb_sub = add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12.33), Inches(0.45))
        tf_sub = tb_sub.text_frame
        p2 = tf_sub.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.add_run()
        run2.text = data["subtitle"]
        run2.font.size = Pt(16)
        run2.font.bold = False
        run2.font.color.rgb = ACCENT

    add_accent_bar(slide)

    # Bullets
    if data.get("bullets"):
        tb_body = add_textbox(slide, Inches(0.55), Inches(1.95), Inches(12.2), Inches(5.25))
        tf_body = tb_body.text_frame
        tf_body.word_wrap = True

        first = True
        for item in data["bullets"]:
            if first:
                p = tf_body.paragraphs[0]
                first = False
            else:
                p = tf_body.add_paragraph()

            if isinstance(item, tuple):
                level, text = item
            else:
                level, text = 0, item

            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(1)
            p.space_after  = Pt(1)

            if text == "":
                p.add_run().text = ""
                continue

            if level == 0:
                # Bold key text if it ends with ':'
                if text.endswith(':'):
                    run = p.add_run()
                    run.text = text
                    run.font.size = Pt(14)
                    run.font.bold = True
                    run.font.color.rgb = GOLD
                else:
                    run = p.add_run()
                    run.text = "• " + text
                    run.font.size = Pt(14)
                    run.font.bold = False
                    run.font.color.rgb = WHITE
            else:
                run = p.add_run()
                run.text = "    ◦ " + text
                run.font.size = Pt(13)
                run.font.bold = False
                run.font.color.rgb = LIGHT_GREY

    return slide


# ── Build ────────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

for i, data in enumerate(SLIDES):
    if data.get("is_title_slide"):
        build_title_slide(prs, data)
    else:
        build_content_slide(prs, data)

os.makedirs(os.path.dirname(os.path.abspath(OUT_PATH)), exist_ok=True)
prs.save(OUT_PATH)
print(f"Saved {len(SLIDES)} slides -> {os.path.abspath(OUT_PATH)}")
